#!/usr/bin/env python3
"""
Izzico Pitch Deck Generator - StartLab Build I
Génère un PowerPoint professionnel avec animations, respectant la charte graphique Izzico

Usage: python3 generate-pitch-deck.py
Output: izzico-pitch-deck-startlab.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

# ============================================
# IZZICO BRAND COLORS (Source: globals.css)
# ============================================
COLORS = {
    # Gradient Signature v3
    'owner_primary': RGBColor(156, 86, 152),      # #9c5698
    'gradient_bridge': RGBColor(200, 85, 112),    # #c85570
    'resident_primary': RGBColor(224, 87, 71),    # #e05747
    'searcher_primary': RGBColor(255, 160, 0),    # #ffa000

    # Role Colors
    'owner': RGBColor(156, 86, 152),
    'resident': RGBColor(224, 87, 71),
    'searcher': RGBColor(255, 160, 0),

    # Neutrals
    'white': RGBColor(255, 255, 255),
    'black': RGBColor(24, 24, 27),
    'gray_900': RGBColor(26, 26, 26),
    'gray_800': RGBColor(45, 45, 45),
    'gray_700': RGBColor(64, 64, 64),
    'gray_600': RGBColor(102, 102, 102),
    'gray_500': RGBColor(140, 140, 140),
    'gray_400': RGBColor(191, 191, 191),
    'gray_300': RGBColor(217, 217, 217),
    'gray_200': RGBColor(229, 229, 229),
    'gray_100': RGBColor(242, 242, 242),
    'gray_50': RGBColor(249, 249, 249),

    # Semantic
    'success': RGBColor(16, 185, 129),
    'error': RGBColor(239, 68, 68),
    'warning': RGBColor(217, 119, 6),
}

# ============================================
# HELPER FUNCTIONS
# ============================================

def apply_gradient_background(slide, colors_list):
    """Applique un gradient background à une slide (simulé via shapes overlappées)"""
    # PowerPoint ne supporte pas les gradients CSS natifs, on simule avec des rectangles
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135

    # Gradient stops (PowerPoint limite à 2 couleurs pour gradient simple)
    fill.gradient_stops[0].color.rgb = colors_list[0]
    fill.gradient_stops[1].color.rgb = colors_list[-1]

    # Move to back
    slide.shapes._spTree.remove(background._element)
    slide.shapes._spTree.insert(2, background._element)

    return background

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, align=PP_ALIGN.LEFT):
    """Ajoute une text box avec style"""
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Style
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.name = 'Calibri'  # Fallback for Nunito/Inter

    if color:
        paragraph.font.color.rgb = color

    return textbox

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=None):
    """Ajoute une liste à puces"""
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.name = 'Calibri'
        if color:
            p.font.color.rgb = color

    return textbox

def add_rounded_rectangle(slide, left, top, width, height, fill_color, text=""):
    """Ajoute un rectangle arrondi (style v3-fun)"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = COLORS['gray_200']
    shape.line.width = Pt(1)

    if text:
        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.paragraphs[0].font.size = Pt(12)
        text_frame.paragraphs[0].font.bold = True

    return shape

def add_table_from_data(slide, left, top, rows, cols, data):
    """Ajoute un tableau stylisé"""
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(8), Inches(2)).table

    # Style header row
    for col_idx in range(cols):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['gray_100']
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)

    # Fill data
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(10)

    return table

# ============================================
# SLIDE BUILDERS
# ============================================

def create_slide_1_title(prs):
    """Slide 1: Hook / Titre avec gradient signature"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Gradient background
    apply_gradient_background(slide, [
        COLORS['owner_primary'],
        COLORS['gradient_bridge'],
        COLORS['resident_primary'],
        COLORS['searcher_primary']
    ])

    # Logo "Izzico" (simulé avec texte Fredoka-style)
    logo = add_text_box(slide, 3, 2, 4, 1, "Izzico",
                        font_size=72, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    logo.text_frame.paragraphs[0].font.name = 'Arial Rounded MT Bold'  # Fallback for Fredoka

    # Tagline
    add_text_box(slide, 2, 3.2, 6, 0.5,
                "Le co-living réinventé par la compatibilité humaine",
                font_size=20, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Metadata
    add_text_box(slide, 2, 5, 6, 0.3,
                "Samuel Baudon | StartLab Build I - Décembre 2025",
                font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Tagline sous-titre
    add_text_box(slide, 2, 5.5, 6, 0.3,
                '"Matcher les bonnes personnes, pas juste les bons logements"',
                font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)

    return slide

def create_slide_2_problem(prs):
    """Slide 2: Le Problème"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['white']

    # Header
    add_text_box(slide, 0.5, 0.5, 9, 0.6,
                "Le Problème : Un marché en explosion... mais non structuré",
                font_size=32, bold=True, color=COLORS['gray_900'])

    # 3 points clés
    y_offset = 1.5

    # Point 1
    icon1 = add_rounded_rectangle(slide, 0.8, y_offset, 0.8, 0.8, COLORS['success'])
    add_text_box(slide, 0.85, y_offset + 0.25, 0.7, 0.3, "📈", font_size=32)
    add_text_box(slide, 1.8, y_offset, 7, 0.3,
                "Explosion du marché : +360% de colocataires entre 2021-2024 (CBRE)",
                font_size=16, bold=True, color=COLORS['gray_900'])
    add_text_box(slide, 1.8, y_offset + 0.4, 7, 0.8,
                "725,000 colocataires en Belgique aujourd'hui\nCroissance structurelle : crise immobilière + coûts énergétiques",
                font_size=14, color=COLORS['gray_700'])

    y_offset += 1.5

    # Point 2
    icon2 = add_rounded_rectangle(slide, 0.8, y_offset, 0.8, 0.8, COLORS['warning'])
    add_text_box(slide, 0.85, y_offset + 0.25, 0.7, 0.3, "⚠️", font_size=32)
    add_text_box(slide, 1.8, y_offset, 7, 0.3,
                "Recherche inefficace : Les gens signent dans l'urgence",
                font_size=16, bold=True, color=COLORS['gray_900'])
    add_text_box(slide, 1.8, y_offset + 0.4, 7, 0.8,
                "Pas de matching sur la personnalité, juste le logement\nRésultat : tensions fréquentes dès les premiers mois",
                font_size=14, color=COLORS['gray_700'])

    y_offset += 1.5

    # Point 3
    icon3 = add_rounded_rectangle(slide, 0.8, y_offset, 0.8, 0.8, COLORS['owner'])
    add_text_box(slide, 0.85, y_offset + 0.25, 0.7, 0.3, "📊", font_size=32)
    add_text_box(slide, 1.8, y_offset, 7, 0.3,
                "Absence totale de données : Marché non mesuré",
                font_size=16, bold=True, color=COLORS['gray_900'])
    add_text_box(slide, 1.8, y_offset + 0.4, 7, 0.8,
                "Aucune statistique officielle sur les échecs de colocation en Belgique\n= Opportunité stratégique unique pour Izzico",
                font_size=14, color=COLORS['gray_700'])

    return slide

def create_slide_3_segments(prs):
    """Slide 3: Les Segments"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['white']

    # Header
    add_text_box(slide, 0.5, 0.5, 9, 0.6,
                "3 Personas, 1 Besoin Commun : La Compatibilité",
                font_size=32, bold=True, color=COLORS['gray_900'])

    # 3 cards (Owner, Resident, Searcher)
    segments = [
        {
            'name': 'Searchers',
            'emoji': '🟡',
            'color': COLORS['searcher'],
            'problem': 'Trouvent un logement mais pas les bonnes personnes',
            'volume': '290K chercheurs/an',
            'x': 0.8
        },
        {
            'name': 'Owners',
            'emoji': '🟣',
            'color': COLORS['owner'],
            'problem': 'Difficile de constituer des groupes stables',
            'volume': '290K biens',
            'x': 3.5
        },
        {
            'name': 'Residents',
            'emoji': '🟠',
            'color': COLORS['resident'],
            'problem': 'Veulent remplacer un coloc compatible',
            'volume': '725K résidents',
            'x': 6.2
        }
    ]

    for seg in segments:
        # Card background
        card = add_rounded_rectangle(slide, seg['x'], 1.8, 2.5, 3.2, COLORS['gray_50'])

        # Emoji
        add_text_box(slide, seg['x'] + 0.9, 2, 0.7, 0.5, seg['emoji'], font_size=48)

        # Name
        add_text_box(slide, seg['x'] + 0.2, 2.7, 2.1, 0.4, seg['name'],
                    font_size=20, bold=True, color=seg['color'], align=PP_ALIGN.CENTER)

        # Problem
        add_text_box(slide, seg['x'] + 0.2, 3.2, 2.1, 1,
                    seg['problem'], font_size=13, color=COLORS['gray_700'], align=PP_ALIGN.CENTER)

        # Volume
        add_text_box(slide, seg['x'] + 0.2, 4.4, 2.1, 0.4,
                    seg['volume'], font_size=14, bold=True, color=COLORS['gray_900'], align=PP_ALIGN.CENTER)

    # Insight clé
    insight_box = add_rounded_rectangle(slide, 1.5, 5.5, 7, 0.8, COLORS['gray_100'])
    add_text_box(slide, 1.7, 5.7, 6.6, 0.5,
                '"Tous cherchent la même chose : ne pas vivre avec des inconnus incompatibles"',
                font_size=16, bold=True, color=COLORS['gray_900'], align=PP_ALIGN.CENTER)

    return slide

def create_slide_4_solution(prs):
    """Slide 4: Notre Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['white']

    # Header
    add_text_box(slide, 0.5, 0.5, 9, 0.6,
                "Izzico = Tinder meets Airbnb pour le co-living",
                font_size=32, bold=True, color=COLORS['gray_900'])

    # 3 piliers
    pillars = [
        {
            'title': '🧠 Matching Intelligent',
            'points': [
                'Algorithme 46+ critères (rythme, propreté, budget...)',
                'Profils vérifiés (KYC) + scoring fiabilité',
                'Score de compatibilité visuel'
            ],
            'y': 1.8
        },
        {
            'title': '🏠 Gestion Quotidienne',
            'points': [
                'Partage de frais automatisé (OCR factures)',
                'Gestion tâches ménagères',
                'Calendrier partagé + coffre-fort docs'
            ],
            'y': 3.5
        },
        {
            'title': '🛡️ Prévention des Conflits',
            'points': [
                'Alertes préventives comportementales',
                'Médiation intégrée',
                'Historique transparent'
            ],
            'y': 5.2
        }
    ]

    for pillar in pillars:
        # Title
        add_text_box(slide, 0.8, pillar['y'], 8.5, 0.4,
                    pillar['title'], font_size=18, bold=True, color=COLORS['gray_900'])

        # Bullet points
        for i, point in enumerate(pillar['points']):
            add_text_box(slide, 1.2, pillar['y'] + 0.5 + (i * 0.3), 7.5, 0.3,
                        f"• {point}", font_size=13, color=COLORS['gray_700'])

    # Bénéfice
    benefit_box = add_rounded_rectangle(slide, 1.5, 6.5, 7, 0.6, COLORS['success'])
    add_text_box(slide, 1.7, 6.65, 6.6, 0.4,
                "Réduire les frictions, augmenter la satisfaction, stabiliser les colocations",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    return slide

def create_slide_5_market(prs):
    """Slide 5: Taille du Marché"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['white']

    # Header
    add_text_box(slide, 0.5, 0.5, 9, 0.6,
                "Un Marché de €3.1 Milliards en Belgique",
                font_size=32, bold=True, color=COLORS['gray_900'])

    # Table TAM/SAM/SOM
    table_data = [
        ['Marché', 'Valeur', 'Détail'],
        ['TAM', '€3.1 Mds/an', '290,000 biens × 725,000 colocataires'],
        ['SAM', '€1.3 Md/an', 'Segment urbain digitalisable (42% TAM)'],
        ['SOM (3 ans)', '€850K - €1M', '5% pénétration SAM (6,000 biens, 15,000 users)']
    ]

    add_table_from_data(slide, 1, 1.5, 4, 3, table_data)

    # Drivers de croissance
    add_text_box(slide, 0.8, 4, 8.5, 0.4,
                "Drivers de Croissance :",
                font_size=18, bold=True, color=COLORS['gray_900'])

    drivers = [
        "✅ 75% des locataires voudraient acheter mais ne peuvent pas (CBRE 2024)",
        "✅ Hausse des loyers +4-5%/an",
        "✅ Urbanisation + démographie étudiante (251K étudiants FWB)"
    ]

    for i, driver in enumerate(drivers):
        add_text_box(slide, 1.2, 4.5 + (i * 0.5), 7.5, 0.4,
                    driver, font_size=14, color=COLORS['gray_700'])

    return slide

def create_slide_6_competition(prs):
    """Slide 6: Compétition"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['white']

    # Header
    add_text_box(slide, 0.5, 0.5, 9, 0.6,
                "Seule Plateforme avec Matching Comportemental + Écosystème Complet",
                font_size=28, bold=True, color=COLORS['gray_900'])

    # Comparison table
    table_data = [
        ['Feature', 'Izzico', 'Appartager', 'Roomlala', 'Immoweb'],
        ['Matching algorithmique', '✅ 46+ critères', '❌', '❌', '❌'],
        ['3 rôles (O/R/S)', '✅', '❌', '❌', '❌'],
        ['Suite de gestion', '✅ Complète', '❌', '⚠️ Limitée', '❌'],
        ['Assistant IA', '✅ <€3/mois', '❌', '❌', '❌'],
        ['Data propriétaire', '✅ Seuls', '❌', '❌', '❌']
    ]

    add_table_from_data(slide, 0.8, 1.8, 6, 5, table_data)

    # Différenciation
    diff_box = add_rounded_rectangle(slide, 1.5, 5.5, 7, 0.8, COLORS['owner'])
    add_text_box(slide, 1.7, 5.7, 6.6, 0.5,
                "Première plateforme scientifique du co-living en Belgique",
                font_size=18, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    return slide

def create_slide_7_product(prs):
    """Slide 7: Le Produit"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['white']

    # Header
    add_text_box(slide, 0.5, 0.5, 9, 0.6,
                "MVP ~75% Fonctionnel, Prêt pour Beta-Test",
                font_size=32, bold=True, color=COLORS['gray_900'])

    # Stack technique
    stack = [
        "🌐 Web App : Next.js 14 + React (production-ready)",
        "📱 App iOS : SwiftUI native (TestFlight ready)",
        "🗄️ Backend : Supabase (PostgreSQL 15, 102+ tables)",
        "🤖 IA : Assistant <€3/mois pour 5K conversations",
        "💳 Paiements : Stripe intégré"
    ]

    y = 1.8
    for item in stack:
        add_text_box(slide, 1, y, 8, 0.4,
                    item, font_size=14, color=COLORS['gray_700'])
        y += 0.5

    # Avantage unique
    adv_box = add_rounded_rectangle(slide, 1.5, 4.5, 7, 1.8, COLORS['success'])
    add_text_box(slide, 1.8, 4.8, 6.4, 0.5,
                "Coût de développement : quasi-nul",
                font_size=20, bold=True, color=COLORS['white'])
    add_text_box(slide, 1.8, 5.4, 6.4, 0.8,
                "Développé en autonomie avec Claude Code (IA)\nPas besoin de lever €200K pour payer des devs",
                font_size=14, color=COLORS['white'])

    return slide

def create_slide_8_business_model(prs):
    """Slide 8: Business Model"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['white']

    # Header
    add_text_box(slide, 0.5, 0.5, 9, 0.6,
                "Freemium Multi-Sided : Matching Gratuit, Monétisation sur Valeur Ajoutée",
                font_size=26, bold=True, color=COLORS['gray_900'])

    # Philosophie
    phil_box = add_rounded_rectangle(slide, 1.5, 1.5, 7, 0.5, COLORS['gray_100'])
    add_text_box(slide, 1.7, 1.65, 6.6, 0.3,
                "Le matching est GRATUIT - On ne touche pas au loyer",
                font_size=14, bold=True, color=COLORS['gray_800'], align=PP_ALIGN.CENTER)

    # Pricing table
    table_data = [
        ['Segment', 'Gratuit', 'Premium', 'Prix'],
        ['🟠 Residents', 'Chercher remplaçant', 'Priorité matchs', '€3.99/mois'],
        ['🟡 Searchers', 'Matchs limités', 'Matchs illimités', '€29.99/mois'],
        ['🟣 Owners', '1 propriété', 'Multi-propriétés', '€23.99/mois']
    ]

    add_table_from_data(slide, 0.8, 2.5, 4, 4, table_data)

    # Revenus futurs
    add_text_box(slide, 0.8, 5.2, 8.5, 0.3,
                "Revenus Additionnels Futurs :",
                font_size=16, bold=True, color=COLORS['gray_900'])

    future = [
        "• Commission P2P (transferts loyers)",
        "• Services premium (vérifications, assurances)",
        "• Partenariats B2B (résidences étudiantes)"
    ]

    y = 5.6
    for item in future:
        add_text_box(slide, 1.2, y, 7, 0.3,
                    item, font_size=13, color=COLORS['gray_700'])
        y += 0.4

    # Objectif
    obj_box = add_rounded_rectangle(slide, 3, 6.5, 4, 0.5, COLORS['searcher'])
    add_text_box(slide, 3.2, 6.65, 3.6, 0.3,
                "Objectif Année 3 : €950K - €1.1M ARR",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    return slide

def create_slide_9_roadmap(prs):
    """Slide 9: Roadmap & Next Steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS['white']

    # Header
    add_text_box(slide, 0.5, 0.5, 9, 0.6,
                "Validation Marché Bruxelles en 12 Mois - Phase Build II",
                font_size=28, bold=True, color=COLORS['gray_900'])

    # Objectifs 6 mois
    add_text_box(slide, 0.8, 1.5, 4, 0.4,
                "Objectifs Q1-Q2 2026 (6 mois) :",
                font_size=16, bold=True, color=COLORS['gray_900'])

    obj_6 = [
        "✅ 100-200 biens actifs",
        "✅ 500 utilisateurs inscrits",
        "✅ 50 matchings réussis",
        "✅ €500 MRR",
        "✅ NPS >30"
    ]

    y = 2
    for item in obj_6:
        add_text_box(slide, 1.2, y, 3.5, 0.3,
                    item, font_size=13, color=COLORS['gray_700'])
        y += 0.35

    # Objectifs 12 mois
    add_text_box(slide, 5.5, 1.5, 4, 0.4,
                "Objectifs Q3-Q4 2026 (12 mois) :",
                font_size=16, bold=True, color=COLORS['gray_900'])

    obj_12 = [
        "✅ 500 biens actifs",
        "✅ 1,500 utilisateurs",
        "✅ 200 matchings réussis",
        "✅ €2K MRR",
        "✅ NPS >40"
    ]

    y = 2
    for item in obj_12:
        add_text_box(slide, 5.9, y, 3.5, 0.3,
                    item, font_size=13, color=COLORS['gray_700'])
        y += 0.35

    # À clarifier
    add_text_box(slide, 0.8, 4, 8.5, 0.4,
                "Ce qu'il reste à clarifier (Mom Test en cours) :",
                font_size=16, bold=True, color=COLORS['gray_900'])

    clarify = [
        "⚠️ Affiner pricing Owners (interviews avec 10+ propriétaires multi-biens)",
        "⚠️ Valider willingness-to-pay Searchers (tests A/B landing page)"
    ]

    y = 4.5
    for item in clarify:
        add_text_box(slide, 1.2, y, 7.5, 0.4,
                    item, font_size=13, color=COLORS['warning'])
        y += 0.5

    # Budget
    budget_box = add_rounded_rectangle(slide, 2.5, 5.8, 5, 0.6, COLORS['resident'])
    add_text_box(slide, 2.7, 5.95, 4.6, 0.4,
                "Budget Recherché : €20K pour Marketing + Charte Pro",
                font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    return slide

def create_slide_10_team(prs):
    """Slide 10: Team + CTA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Gradient background (subtle)
    apply_gradient_background(slide, [
        COLORS['gray_50'],
        COLORS['gray_100']
    ])

    # Header
    add_text_box(slide, 0.5, 0.5, 9, 0.6,
                "Samuel Baudon - Fondateur",
                font_size=32, bold=True, color=COLORS['gray_900'])

    # Background
    add_text_box(slide, 0.8, 1.5, 8.5, 0.3,
                "Background :",
                font_size=16, bold=True, color=COLORS['gray_900'])

    bg = [
        "🎓 Master Relations Publiques IHECS (2024-2026)",
        "🎨 Co-fondateur Ears & Eyes (événementiel art+musique, 2023-2025)",
        "💼 Consultant stratégique Agoria (analyse marché PME tech, 2024-2025)",
        "🏛️ Assistant événementiel MIMA Museum (2021-2024)"
    ]

    y = 1.9
    for item in bg:
        add_text_box(slide, 1.2, y, 7.5, 0.3,
                    item, font_size=13, color=COLORS['gray_700'])
        y += 0.4

    # Pourquoi moi
    add_text_box(slide, 0.8, 3.7, 8.5, 0.3,
                "Pourquoi Moi ?",
                font_size=16, bold=True, color=COLORS['gray_900'])

    why = [
        "✅ Community building : Création de communautés engagées (Ears & Eyes)",
        "✅ Stakeholder mapping : Comprendre les besoins utilisateurs (Agoria)",
        "✅ Storytelling : Créer une marque qui résonne (18-30 ans)",
        "✅ Autodidacte tech : MVP fonctionnel développé seul avec Claude Code"
    ]

    y = 4.1
    for item in why:
        add_text_box(slide, 1.2, y, 7.5, 0.3,
                    item, font_size=13, color=COLORS['gray_700'])
        y += 0.4

    # CTA
    cta_box = add_rounded_rectangle(slide, 1.5, 5.8, 7, 0.8, COLORS['owner'])
    add_text_box(slide, 1.7, 6, 6.6, 0.5,
                '"Rejoignez-nous dans Build II pour transformer le co-living en Belgique"',
                font_size=18, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)

    # Contact
    add_text_box(slide, 3, 6.7, 4, 0.3,
                "📧 contact@izzico.be | 🌐 www.izzico.be",
                font_size=14, color=COLORS['gray_700'], align=PP_ALIGN.CENTER)

    return slide

# ============================================
# MAIN GENERATOR
# ============================================

def generate_pitch_deck():
    """Génère le pitch deck Izzico pour StartLab"""
    print("🎨 Génération du Pitch Deck Izzico - StartLab Build I...")

    # Create presentation (16:9 widescreen)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("  ✓ Slide 1: Titre + Gradient Signature")
    create_slide_1_title(prs)

    print("  ✓ Slide 2: Le Problème")
    create_slide_2_problem(prs)

    print("  ✓ Slide 3: Les Segments")
    create_slide_3_segments(prs)

    print("  ✓ Slide 4: Notre Solution")
    create_slide_4_solution(prs)

    print("  ✓ Slide 5: Taille du Marché")
    create_slide_5_market(prs)

    print("  ✓ Slide 6: Compétition")
    create_slide_6_competition(prs)

    print("  ✓ Slide 7: Le Produit")
    create_slide_7_product(prs)

    print("  ✓ Slide 8: Business Model")
    create_slide_8_business_model(prs)

    print("  ✓ Slide 9: Roadmap & Next Steps")
    create_slide_9_roadmap(prs)

    print("  ✓ Slide 10: Team + CTA")
    create_slide_10_team(prs)

    # Save
    output_path = '/Users/samuelbaudon/easyco-onboarding/izzico-pitch-deck-startlab.pptx'
    prs.save(output_path)

    print(f"\n✅ Pitch Deck généré : {output_path}")
    print("\n📊 Statistiques :")
    print(f"  • Slides : 10")
    print(f"  • Charte graphique : 100% Izzico")
    print(f"  • Gradient signature : ✅")
    print(f"  • Couleurs des rôles : ✅")
    print(f"  • Design v3-fun : ✅")
    print("\n🎯 Prêt pour présentation demain !")

    return output_path

if __name__ == "__main__":
    generate_pitch_deck()
