#!/usr/bin/env python3
"""
Generate Izzico Pitch Deck as PowerPoint (.pptx)
Reproduit EXACTEMENT le design du HTML scrollable
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from PIL import Image
import qrcode
import io
import os

# Izzico Brand Colors (exactement comme dans le HTML)
OWNER = RGBColor(156, 86, 152)          # #9c5698
OWNER_LIGHT = RGBColor(248, 240, 247)   # #F8F0F7
OWNER_HOVER = RGBColor(176, 112, 168)   # #B070A8

RESIDENT = RGBColor(224, 87, 71)        # #e05747
RESIDENT_LIGHT = RGBColor(254, 242, 238) # #FEF2EE
RESIDENT_HOVER = RGBColor(233, 106, 80)  # #E96A50

SEARCHER = RGBColor(255, 160, 0)        # #ffa000
SEARCHER_LIGHT = RGBColor(255, 251, 235) # #FFFBEB
SEARCHER_HOVER = RGBColor(251, 191, 36)  # #FBBF24

SUCCESS = RGBColor(124, 184, 155)       # #7CB89B
SUCCESS_LIGHT = RGBColor(240, 247, 244) # #F0F7F4
WARNING = RGBColor(217, 168, 112)       # #D9A870
WARNING_LIGHT = RGBColor(255, 251, 235) # #FFFBEB

GRAY_900 = RGBColor(26, 26, 26)
GRAY_700 = RGBColor(64, 64, 64)
GRAY_600 = RGBColor(102, 102, 102)
GRAY_200 = RGBColor(229, 229, 229)
GRAY_100 = RGBColor(242, 242, 242)
GRAY_50 = RGBColor(249, 249, 249)
WHITE = RGBColor(255, 255, 255)

# Gradient signature colors
GRADIENT_BRIDGE = RGBColor(200, 85, 112)  # #c85570

def add_gradient_background(slide, prs, colors):
    """Add gradient background à une slide"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0,
        prs.slide_width, prs.slide_height
    )

    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 135.0

    # Python-pptx gradient API: use index-based access
    gradient_stops = fill.gradient_stops

    # Clear existing stops and set new ones
    # First stop
    gradient_stops[0].color.rgb = colors[0][1]
    gradient_stops[0].position = colors[0][0]

    # Last stop
    gradient_stops[1].color.rgb = colors[-1][1]
    gradient_stops[1].position = colors[-1][0]

    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, content, color_scheme='neutral', icon_text=''):
    """Add une card colorée comme dans le HTML"""
    # Card background
    card = slide.shapes.add_shape(1, left, top, width, height)

    # Apply color scheme
    if color_scheme == 'success':
        card.fill.solid()
        card.fill.fore_color.rgb = SUCCESS_LIGHT
        card.line.color.rgb = SUCCESS
        border_color = SUCCESS
    elif color_scheme == 'warning':
        card.fill.solid()
        card.fill.fore_color.rgb = WARNING_LIGHT
        card.line.color.rgb = WARNING
        border_color = WARNING
    elif color_scheme == 'owner':
        card.fill.solid()
        card.fill.fore_color.rgb = OWNER_LIGHT
        card.line.color.rgb = OWNER
        border_color = OWNER
    elif color_scheme == 'resident':
        card.fill.solid()
        card.fill.fore_color.rgb = RESIDENT_LIGHT
        card.line.color.rgb = RESIDENT
        border_color = RESIDENT
    elif color_scheme == 'searcher':
        card.fill.solid()
        card.fill.fore_color.rgb = SEARCHER_LIGHT
        card.line.color.rgb = SEARCHER
        border_color = SEARCHER
    else:
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = GRAY_200
        border_color = GRAY_200

    card.line.width = Pt(3)

    # Add icon circle (simulated)
    icon_size = Inches(0.5)
    icon = slide.shapes.add_shape(
        1,  # Rectangle/Circle
        left + Inches(0.15), top + Inches(0.15),
        icon_size, icon_size
    )
    icon.fill.solid()
    icon.fill.fore_color.rgb = border_color
    icon.line.fill.background()

    # Add text
    text_box = slide.shapes.add_textbox(
        left + Inches(0.8), top + Inches(0.15),
        width - Inches(0.95), height - Inches(0.3)
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # Title
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GRAY_900
    p.space_after = Pt(8)

    # Content
    if isinstance(content, list):
        for item in content:
            p = tf.add_paragraph()
            p.text = f"  • {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = GRAY_700
            p.level = 0
    else:
        p = tf.add_paragraph()
        p.text = content
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY_700

def add_slide_header(slide, prs, text):
    """Add header avec gradient border"""
    # Header text
    header = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.7)
    )
    tf = header.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GRAY_900

    # Gradient line below (simulated with colored rectangle)
    line = slide.shapes.add_shape(
        1,
        Inches(0.5), Inches(1.05),
        Inches(9), Inches(0.05)
    )
    # Gradient fill (simplified to 2 colors)
    fill = line.fill
    fill.gradient()
    fill.gradient_angle = 90.0

    fill.gradient_stops[0].color.rgb = OWNER
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = SEARCHER
    fill.gradient_stops[1].position = 1.0

    line.line.fill.background()

def generate_qr_code(data, size=300):
    """Generate QR code as PIL Image"""
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_H,
        box_size=10,
        border=2,
    )
    qr.add_data(data)
    qr.make(fit=True)

    img = qr.make_image(fill_color="black", back_color="white")
    return img

# ============================================
# SLIDE GENERATORS
# ============================================

def slide_1_title(prs):
    """Slide 1: Titre avec gradient signature"""
    print("  📄 Slide 1: Titre...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Gradient background
    add_gradient_background(slide, prs, [
        (0.0, OWNER),
        (0.20, GRADIENT_BRIDGE),
        (0.35, RGBColor(209, 86, 89)),
        (0.50, RESIDENT),
        (0.75, RGBColor(255, 124, 16)),
        (1.0, SEARCHER)
    ])

    # Logo placeholder (simplified)
    logo_box = slide.shapes.add_textbox(Inches(2.5), Inches(1.5), Inches(5), Inches(1))
    tf = logo_box.text_frame
    tf.text = "izzico"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(90)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Fredoka'

    # Tagline
    tagline = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    tf = tagline.text_frame
    tf.text = "Le co-living réinventé par la compatibilité humaine"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Nunito'

    # Metadata
    meta = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.5))
    tf = meta.text_frame
    tf.text = "Samuel Baudon | StartLab Build I - Décembre 2025"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.name = 'Inter'

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(0.5))
    tf = subtitle.text_frame
    tf.text = '"Matcher les bonnes personnes, pas juste les bons logements"'
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(26)
    p.font.italic = True
    p.font.color.rgb = WHITE
    p.font.name = 'Inter'

def slide_2_problem(prs):
    """Slide 2: Le Problème"""
    print("  📄 Slide 2: Le Problème...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RESIDENT_LIGHT
    bg.line.fill.background()

    add_slide_header(slide, prs, "Le Problème : Un Marché en Explosion... Mais Non Structuré")

    # Card 1: Explosion
    add_card(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(1.4),
             "Explosion du Marché",
             ["+360% de colocataires entre 2021-2024 (CBRE)",
              "725,000 colocataires en Belgique aujourd'hui"],
             'success')

    # Card 2: Recherche Inefficace
    add_card(slide, Inches(0.5), Inches(2.9), Inches(9), Inches(1.4),
             "Recherche Inefficace",
             ["Les gens signent dans l'urgence sans connaître leurs futurs colocs",
              "Pas de matching sur la personnalité"],
             'warning')

    # Card 3: Absence de Données
    add_card(slide, Inches(0.5), Inches(4.5), Inches(9), Inches(1.4),
             "Absence Totale de Données",
             ["Marché non mesuré, non structuré",
              "= Opportunité stratégique unique pour Izzico"],
             'owner')

def slide_3_personas(prs):
    """Slide 3: 3 Personas"""
    print("  📄 Slide 3: 3 Personas...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SEARCHER_LIGHT
    bg.line.fill.background()

    add_slide_header(slide, prs, "3 Personas, 1 Besoin Commun : La Compatibilité")

    # Personas grid (3 cards)
    personas = [
        ("Searchers", "Trouvent un logement mais pas les bonnes personnes",
         "Volume = turnover (?)", SEARCHER, SEARCHER_LIGHT),
        ("Owners", "Difficile de constituer des groupes stables",
         "290K biens", OWNER, OWNER_LIGHT),
        ("Residents", "Veulent remplacer un coloc compatible",
         "725K résidents", RESIDENT, RESIDENT_LIGHT)
    ]

    card_width = Inches(2.8)
    start_x = Inches(0.6)
    gap = Inches(0.3)

    for i, (name, problem, volume, color, bg_color) in enumerate(personas):
        x = start_x + i * (card_width + gap)

        # Card
        card = slide.shapes.add_shape(1, x, Inches(1.4), card_width, Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = color
        card.line.width = Pt(4)

        # Icon circle
        icon = slide.shapes.add_shape(
            1, x + Inches(0.9), Inches(1.6), Inches(1), Inches(1)
        )
        icon.fill.solid()
        icon.fill.fore_color.rgb = color
        icon.line.fill.background()

        # Name
        name_box = slide.shapes.add_textbox(x, Inches(2.8), card_width, Inches(0.5))
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = GRAY_900
        p.font.name = 'Nunito'

        # Problem
        prob_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.3), card_width - Inches(0.4), Inches(0.8))
        tf = prob_box.text_frame
        tf.text = problem
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY_700

        # Volume badge
        vol_box = slide.shapes.add_textbox(x + Inches(0.4), Inches(4.2), card_width - Inches(0.8), Inches(0.4))
        vol_box.fill.solid()
        vol_box.fill.fore_color.rgb = color
        tf = vol_box.text_frame
        tf.text = volume
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # Highlight box
    highlight = slide.shapes.add_shape(1, Inches(0.6), Inches(4.8), Inches(8.8), Inches(0.6))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = GRAY_100
    highlight.line.fill.background()

    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.9), Inches(8.4), Inches(0.4))
    tf = text_box.text_frame
    tf.text = '"Tous cherchent la même chose : ne pas vivre avec des inconnus incompatibles"'
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GRAY_900

def slide_4_solution(prs):
    """Slide 4: Notre Solution"""
    print("  📄 Slide 4: Solution...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = OWNER_LIGHT
    bg.line.fill.background()

    add_slide_header(slide, prs, "Izzico = Tinder meets Airbnb pour le co-living")

    # 3 feature cards
    add_card(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(1.1),
             "Matching Intelligent",
             ["Algorithme 46+ critères (rythme, propreté, budget...)",
              "Profils vérifiés (KYC) + scoring"],
             'resident')

    add_card(slide, Inches(0.5), Inches(2.7), Inches(9), Inches(1.1),
             "Gestion Quotidienne",
             ["Partage de frais automatisé (OCR factures)",
              "Gestion tâches ménagères + calendrier"],
             'owner')

    add_card(slide, Inches(0.5), Inches(4.0), Inches(9), Inches(1.1),
             "Prévention des Conflits",
             ["Alertes préventives basées sur données",
              "Médiation intégrée + historique"],
             'searcher')

    # Highlight bottom
    highlight = slide.shapes.add_shape(1, Inches(0.5), Inches(5.3), Inches(9), Inches(0.5))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = SUCCESS
    highlight.line.fill.background()

    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.35), Inches(8.6), Inches(0.4))
    tf = text_box.text_frame
    tf.text = "Réduire les frictions, augmenter la satisfaction, stabiliser les colocations"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE

def slide_5_market(prs):
    """Slide 5: Taille du Marché"""
    print("  📄 Slide 5: Marché...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RESIDENT_LIGHT
    bg.line.fill.background()

    add_slide_header(slide, prs, "Un Marché de €3.1 Milliards en Belgique")

    # 3 stat cards
    stats = [
        ("€3.1 Mds", "TAM - Marché Total", "290K biens × 725K colocataires", OWNER),
        ("€1.3 Md", "SAM - Digitalisable", "Zones urbaines (42% TAM)", RESIDENT),
        ("€1M", "SOM - Objectif 3 ans", "5% pénétration SAM", SEARCHER)
    ]

    card_width = Inches(2.8)
    for i, (number, label, detail, color) in enumerate(stats):
        x = Inches(0.6 + i * 3.1)

        card = slide.shapes.add_shape(1, x, Inches(1.5), card_width, Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = color
        card.line.width = Pt(4)

        # Number
        num_box = slide.shapes.add_textbox(x, Inches(1.7), card_width, Inches(0.8))
        tf = num_box.text_frame
        tf.text = number
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = color

        # Label
        label_box = slide.shapes.add_textbox(x, Inches(2.5), card_width, Inches(0.4))
        tf = label_box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GRAY_600

        # Detail
        detail_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.9), card_width - Inches(0.4), Inches(0.5))
        tf = detail_box.text_frame
        tf.text = detail
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY_600

    # Drivers card
    add_card(slide, Inches(0.5), Inches(3.7), Inches(9), Inches(1.5),
             "Drivers de Croissance",
             ["75% des locataires voudraient acheter (CBRE 2024)",
              "Hausse des loyers +4-5%/an",
              "Urbanisation + démographie étudiante"],
             'success')

def slide_13_screenshot_qr(prs):
    """Slide 13: Screenshot + QR Code"""
    print("  📄 Slide 13: Screenshot + QR...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background gris
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = GRAY_100
    bg.line.fill.background()

    # Try to add screenshot if exists
    screenshot_path = os.path.join(os.path.dirname(__file__), '..', 'izzico-homepage-screenshot.png')
    if os.path.exists(screenshot_path):
        # Add screenshot as background
        slide.shapes.add_picture(
            screenshot_path,
            0, 0,
            width=prs.slide_width,
            height=prs.slide_height
        )

    # QR Code box (glassmorphism simulated)
    qr_box = slide.shapes.add_shape(1, Inches(3.5), Inches(1.5), Inches(3), Inches(3))
    qr_box.fill.solid()
    qr_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    qr_box.fill.transparency = 0.85  # 15% opacity
    qr_box.line.color.rgb = WHITE
    qr_box.line.width = Pt(1)

    # Generate and add QR code
    qr_img = generate_qr_code("https://izzico.be", size=300)

    # Save QR to temp file
    qr_path = '/tmp/izzico_qr.png'
    qr_img.save(qr_path)

    # Add QR to slide
    slide.shapes.add_picture(
        qr_path,
        Inches(4), Inches(2),
        width=Inches(2),
        height=Inches(2)
    )

    # URL box
    url_box = slide.shapes.add_shape(1, Inches(3.5), Inches(4.7), Inches(3), Inches(0.6))
    url_box.fill.solid()
    url_box.fill.fore_color.rgb = WHITE
    url_box.fill.transparency = 0.85
    url_box.line.color.rgb = WHITE
    url_box.line.width = Pt(1)

    text = slide.shapes.add_textbox(Inches(3.5), Inches(4.8), Inches(3), Inches(0.4))
    tf = text.text_frame
    tf.text = "izzico.be"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Fredoka'

def slide_14_team(prs):
    """Slide 14: Team + CTA"""
    print("  📄 Slide 14: Team...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Gradient background
    add_gradient_background(slide, prs, [
        (0.0, OWNER),
        (0.50, RESIDENT),
        (1.0, SEARCHER)
    ])

    add_slide_header(slide, prs, "Samuel Baudon - Fondateur")

    # Card 1: Background
    card1 = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(4.2), Inches(2.5))
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card1.fill.transparency = 0.05
    card1.line.color.rgb = WHITE
    card1.line.width = Pt(2)

    text1 = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(3.8), Inches(2.3))
    tf = text1.text_frame
    tf.text = "Background\n\n• Master Relations Publiques IHECS\n• Co-fondateur Ears & Eyes\n• Consultant Agoria\n• Assistant MIMA Museum"
    for i, p in enumerate(tf.paragraphs):
        p.font.size = Pt(18) if i == 0 else Pt(14)
        p.font.bold = (i == 0)
        p.font.color.rgb = WHITE
        p.font.name = 'Inter'

    # Card 2: Pourquoi
    card2 = slide.shapes.add_shape(1, Inches(5.3), Inches(1.4), Inches(4.2), Inches(2.5))
    card2.fill.solid()
    card2.fill.fore_color.rgb = WHITE
    card2.fill.transparency = 0.05
    card2.line.color.rgb = WHITE
    card2.line.width = Pt(2)

    text2 = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(3.8), Inches(2.3))
    tf = text2.text_frame
    tf.text = "Pourquoi Izzico ?\n\n• 18-30 ans : Je suis dans la cible\n• Claude Code : Outil de création\n• Agoria : Professionnalisme\n• Ears & Eyes : Vision communauté"
    for i, p in enumerate(tf.paragraphs):
        p.font.size = Pt(18) if i == 0 else Pt(14)
        p.font.bold = (i == 0)
        p.font.color.rgb = WHITE
        p.font.name = 'Inter'

    # CTA box
    cta = slide.shapes.add_shape(1, Inches(1), Inches(4.2), Inches(8), Inches(1.2))
    cta.fill.solid()
    cta.fill.fore_color.rgb = WHITE
    cta.fill.transparency = 0.75
    cta.line.color.rgb = WHITE
    cta.line.width = Pt(3)

    cta_text = slide.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(7.6), Inches(1))
    tf = cta_text.text_frame
    tf.text = '"Rejoignez-nous pour transformer le co-living"\n\nhello@izzico.be | www.izzico.be'

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Nunito'

    p2 = tf.paragraphs[1]
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = WHITE

def create_full_deck():
    """Generate complete 14-slide deck"""
    print("🚀 Génération du Pitch Deck Izzico (.pptx)...\n")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # Generate all slides
    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_personas(prs)
    slide_4_solution(prs)
    slide_5_market(prs)

    # Slides 6-12: Add placeholders for now
    for i in range(6, 13):
        print(f"  📄 Slide {i}: Placeholder...")
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = GRAY_50
        bg.line.fill.background()

        text = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1))
        tf = text.text_frame
        tf.text = f"Slide {i} - À compléter"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(48)
        p.font.color.rgb = GRAY_700

    slide_13_screenshot_qr(prs)
    slide_14_team(prs)

    # Save
    output_path = os.path.join(os.path.dirname(__file__), '..', 'izzico-pitch-deck-startlab.pptx')
    prs.save(output_path)

    file_size = os.path.getsize(output_path) / (1024 * 1024)

    print(f"\n✅ PowerPoint généré avec succès !")
    print(f"   📦 Fichier : {output_path}")
    print(f"   📊 Taille : {file_size:.2f} MB")
    print(f"   📐 Format : 16:9 (1920×1080 équivalent)")
    print(f"\n📋 Prochaines étapes :")
    print(f"   1. Ouvrir dans PowerPoint/Keynote")
    print(f"   2. Compléter les slides 6-12")
    print(f"   3. Exporter en PDF (Fichier > Exporter > PDF)")
    print(f"   4. Le PDF PowerPoint sera parfait !")

if __name__ == '__main__':
    create_full_deck()
